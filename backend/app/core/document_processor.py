from typing import List, Dict, Any
from pathlib import Path
import logging
try:
    from langchain.text_splitter import RecursiveCharacterTextSplitter
except ImportError:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
try:
    from langchain_openai import OpenAIEmbeddings
except ImportError:
    OpenAIEmbeddings = None

try:
    from langchain_community.embeddings import HuggingFaceEmbeddings
except ImportError:
    HuggingFaceEmbeddings = None
import chromadb
from chromadb.config import Settings as ChromaSettings
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl
from bs4 import BeautifulSoup

from app.core.config import settings

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Advanced document processing with semantic chunking"""
    
    def __init__(self):
        # Use local embeddings if configured
        if getattr(settings, 'USE_LOCAL_MODELS', False):
            if HuggingFaceEmbeddings is None:
                raise ImportError("Please install: pip install sentence-transformers")
            self.embeddings = HuggingFaceEmbeddings(
                model_name=settings.LOCAL_EMBEDDING_MODEL,
                model_kwargs={'device': 'cpu'},
                encode_kwargs={'normalize_embeddings': True}
            )
        else:
            if OpenAIEmbeddings is None:
                raise ImportError("Please install: pip install langchain-openai")
            self.embeddings = OpenAIEmbeddings(
                model="text-embedding-3-small",
                openai_api_key=settings.OPENAI_API_KEY
            )
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=settings.CHUNK_SIZE,
            chunk_overlap=settings.CHUNK_OVERLAP,
            separators=["\n\n", "\n", ". ", " ", ""]
        )
        self.chroma_client = chromadb.PersistentClient(
            path=settings.CHROMA_PERSIST_DIR,
            settings=ChromaSettings(anonymized_telemetry=False)
        )
    
    async def process_document(
        self,
        file_path: str,
        file_type: str,
        tenant_id: int,
        document_id: int,
        metadata: Dict[str, Any] = None
    ) -> Dict[str, Any]:
        """Process document and store in vector database"""
        try:
            # Extract text based on file type
            text = await self.extract_text(file_path, file_type)
            
            # Smart chunking
            chunks = await self.smart_chunking(text, metadata)
            
            # Generate embeddings and store
            chunk_count = await self.store_chunks(
                chunks,
                tenant_id,
                document_id,
                metadata
            )
            
            return {
                "status": "success",
                "chunk_count": chunk_count,
                "text_length": len(text)
            }
            
        except Exception as e:
            logger.error(f"Error processing document: {str(e)}")
            return {
                "status": "failed",
                "error": str(e)
            }
    
    async def extract_text(self, file_path: str, file_type: str) -> str:
        """Extract text from various file formats"""
        
        if file_type == "pdf":
            return await self._extract_pdf(file_path)
        elif file_type in ["docx", "doc"]:
            return await self._extract_docx(file_path)
        elif file_type in ["pptx", "ppt"]:
            return await self._extract_pptx(file_path)
        elif file_type in ["xlsx", "xls"]:
            return await self._extract_excel(file_path)
        elif file_type == "html":
            return await self._extract_html(file_path)
        elif file_type == "txt":
            return await self._extract_txt(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
    
    async def _extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF"""
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n\n"
        return text
    
    async def _extract_docx(self, file_path: str) -> str:
        """Extract text from DOCX"""
        doc = DocxDocument(file_path)
        text = "\n\n".join([paragraph.text for paragraph in doc.paragraphs])
        return text
    
    async def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PPTX"""
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            text += "\n"
        return text
    
    async def _extract_excel(self, file_path: str) -> str:
        """Extract text from Excel"""
        wb = openpyxl.load_workbook(file_path)
        text = ""
        for sheet in wb.worksheets:
            text += f"Sheet: {sheet.title}\n"
            for row in sheet.iter_rows(values_only=True):
                text += " | ".join([str(cell) if cell else "" for cell in row]) + "\n"
            text += "\n"
        return text
    
    async def _extract_html(self, file_path: str) -> str:
        """Extract text from HTML"""
        with open(file_path, 'r', encoding='utf-8') as file:
            soup = BeautifulSoup(file.read(), 'html.parser')
            text = soup.get_text(separator='\n')
        return text
    
    async def _extract_txt(self, file_path: str) -> str:
        """Extract text from TXT"""
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    
    async def smart_chunking(
        self,
        text: str,
        metadata: Dict[str, Any] = None
    ) -> List[Dict[str, Any]]:
        """Smart chunking with semantic boundaries"""
        
        # Split text into chunks
        chunks = self.text_splitter.split_text(text)
        
        # Add metadata to each chunk
        chunk_dicts = []
        for i, chunk in enumerate(chunks):
            chunk_dict = {
                "content": chunk,
                "metadata": {
                    **(metadata or {}),
                    "chunk_index": i,
                    "total_chunks": len(chunks)
                }
            }
            chunk_dicts.append(chunk_dict)
        
        return chunk_dicts
    
    async def store_chunks(
        self,
        chunks: List[Dict[str, Any]],
        tenant_id: int,
        document_id: int,
        metadata: Dict[str, Any] = None
    ) -> int:
        """Store chunks in vector database"""
        
        collection_name = f"tenant_{tenant_id}"
        
        try:
            collection = self.chroma_client.get_or_create_collection(
                name=collection_name,
                metadata={"tenant_id": tenant_id}
            )
        except Exception as e:
            logger.error(f"Error creating collection: {str(e)}")
            raise
        
        # Prepare data for ChromaDB
        documents = [chunk["content"] for chunk in chunks]
        metadatas = [
            {
                **chunk["metadata"],
                "document_id": document_id,
                **(metadata or {})
            }
            for chunk in chunks
        ]
        ids = [f"doc_{document_id}_chunk_{i}" for i in range(len(chunks))]
        
        # Generate embeddings
        embeddings = self.embeddings.embed_documents(documents)
        
        # Store in ChromaDB
        collection.add(
            documents=documents,
            embeddings=embeddings,
            metadatas=metadatas,
            ids=ids
        )
        
        return len(chunks)
    
    async def delete_document_chunks(self, tenant_id: int, document_id: int):
        """Delete all chunks for a document"""
        collection_name = f"tenant_{tenant_id}"
        
        try:
            collection = self.chroma_client.get_collection(collection_name)
            # Delete chunks with matching document_id
            collection.delete(where={"document_id": document_id})
        except Exception as e:
            logger.error(f"Error deleting chunks: {str(e)}")
